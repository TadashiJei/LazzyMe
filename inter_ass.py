import eel
import speech_recognition as sr
from openai import OpenAI
import google.generativeai as genai
import threading
import json
import os
import base64
import time
import re
import numpy as np
import queue
import io
from array import array
import wave
from docx import Document
from PyPDF2 import PdfReader
import pandas as pd
from pptx import Presentation
import pptx
import tempfile
from pathlib import Path

eel.init('web')

class AudioAssistant:
    def __init__(self):
        self.setup_audio()
        self.is_listening = False
        self.openai_client = None
        self.gemini_client = None
        self.api_keys = {}
        self.current_model = "gpt"
        self.current_personality = "professional"
        self.current_language = "en-US"
        self.current_role = "interviewer"
        self.knowledge_base = ""
        self.company_info = ""
        self.job_guidelines = ""
        self.hackathon_info = {
            "name": "",
            "problem_statement": "",
            "criteria": [],
            "sdg_goals": [],
            "additional_context": ""
        }
        self.audio_buffer = queue.Queue()
        self.is_processing_audio = False
        self.is_speaking = False
        self.audio_playing = False
        self.is_muted = False
        self.tts_enabled = True
        self.load_config()
        self.temp_dir = tempfile.mkdtemp()

    def setup_audio(self):
        self.recognizer = sr.Recognizer()
        self.mic = sr.Microphone()
        with self.mic as source:
            self.recognizer.adjust_for_ambient_noise(source)
        # Adjust recognition settings
        self.recognizer.energy_threshold = 300  # Increase sensitivity
        self.recognizer.dynamic_energy_threshold = True
        self.recognizer.dynamic_energy_adjustment_damping = 0.15
        self.recognizer.dynamic_energy_ratio = 1.5
        self.recognizer.pause_threshold = 0.8  # Shorter pause threshold
        self.recognizer.phrase_threshold = 0.3
        self.recognizer.non_speaking_duration = 0.5

    def process_audio_data(self, audio_data):
        try:
            # Convert audio data to wav format
            audio_array = array('h', audio_data)
            with io.BytesIO() as wav_buffer:
                with wave.open(wav_buffer, 'wb') as wav_file:
                    wav_file.setnchannels(1)
                    wav_file.setsampwidth(2)
                    wav_file.setframerate(44100)
                    wav_file.writeframes(audio_array.tobytes())
                
                wav_buffer.seek(0)
                with sr.AudioFile(wav_buffer) as source:
                    audio = self.recognizer.record(source)
                    try:
                        # Try multiple recognition services for better accuracy
                        text = None
                        try:
                            # Try Google Cloud Speech first (better with accents)
                            text = self.recognizer.recognize_google_cloud(
                                audio,
                                language=self.current_language,
                                preferred_phrases=["interview", "question", "answer", "experience", "skills"]
                            )
                        except:
                            # Fallback to regular Google Speech Recognition
                            text = self.recognizer.recognize_google(
                                audio,
                                language=self.current_language
                            )
                        
                        if text:
                            if self.is_question(text):
                                capitalized_text = text[0].upper() + text[1:]
                                if not capitalized_text.endswith('?'):
                                    capitalized_text += '?'
                                eel.update_ui(f"Q: {capitalized_text}", "")
                                self.is_speaking = True
                                response = self.get_ai_response(capitalized_text)
                                eel.update_ui("", f"{response}")
                                self.is_speaking = False
                    except sr.UnknownValueError:
                        pass  # No speech detected
                    except sr.RequestError as e:
                        print(f"Could not request results; {e}")
                    except Exception as e:
                        print(f"Error processing audio: {str(e)}")
        except Exception as e:
            print(f"Error processing audio data: {str(e)}")

    def process_uploaded_file(self, file_path):
        try:
            ext = Path(file_path).suffix.lower()
            content = ""
            
            if ext == '.pdf':
                with open(file_path, 'rb') as file:
                    pdf = PdfReader(file)
                    content = "\n".join(page.extract_text() for page in pdf.pages)
            
            elif ext == '.docx':
                doc = Document(file_path)
                content = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            
            elif ext == '.pptx':
                prs = Presentation(file_path)
                content = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            
            elif ext == '.csv':
                df = pd.read_csv(file_path)
                content = df.to_string()
            
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as file:
                    content = file.read()
            
            return content
        except Exception as e:
            print(f"Error processing file {file_path}: {str(e)}")
            return ""

    def load_config(self):
        if os.path.exists('config.json'):
            with open('config.json', 'r') as f:
                config = json.load(f)
                self.api_keys = config.get('api_keys', {})
                self.current_model = config.get('current_model', 'gpt')
                self.current_personality = config.get('current_personality', 'professional')
                self.current_language = config.get('current_language', 'en-US')
                self.current_role = config.get('current_role', 'interviewer')
                self.job_guidelines = config.get('job_guidelines', '')
                self.knowledge_base = config.get('knowledge_base', '')
                self.company_info = config.get('company_info', '')
                self.hackathon_info = config.get('hackathon_info', {
                    "name": "",
                    "problem_statement": "",
                    "criteria": [],
                    "sdg_goals": [],
                    "additional_context": ""
                })
                self.setup_clients()

    def save_config(self):
        with open('config.json', 'w') as f:
            json.dump({
                'api_keys': self.api_keys,
                'current_model': self.current_model,
                'current_personality': self.current_personality,
                'current_language': self.current_language,
                'current_role': self.current_role,
                'job_guidelines': self.job_guidelines,
                'knowledge_base': self.knowledge_base,
                'company_info': self.company_info,
                'hackathon_info': self.hackathon_info
            }, f)

    def set_job_guidelines(self, guidelines):
        self.job_guidelines = guidelines
        self.save_config()

    def get_job_guidelines(self):
        return self.job_guidelines

    def start_screen_audio_capture(self):
        if not self.system_audio_device:
            error_msg = """
            No virtual audio device found. To capture system audio on macOS, you need to:
            1. Install BlackHole (recommended) or Soundflower
            2. Set up an Aggregate Device or Multi-Output Device in Audio MIDI Setup
            3. Configure your system to output audio to the virtual device
            
            Would you like to install BlackHole now? You can download it from:
            https://existential.audio/blackhole/
            """
            print(error_msg)
            return {"success": False, "message": error_msg}
        
        self.is_capturing_screen = True
        threading.Thread(target=self._capture_screen_audio, daemon=True).start()
        return {"success": True, "message": f"Started capturing audio from {self.system_audio_device['name']}"}

    def stop_screen_audio_capture(self):
        self.is_capturing_screen = False

    def _capture_screen_audio(self):
        def audio_callback(indata, frames, time, status):
            if status:
                print(f'Audio callback status: {status}')
            self.screen_audio_queue.put(indata.copy())

        try:
            with sd.InputStream(device=self.system_audio_device['index'],
                              channels=min(2, self.system_audio_device['max_input_channels']),
                              callback=audio_callback,
                              blocksize=8192,
                              samplerate=44100):
                print(f"Started audio capture from {self.system_audio_device['name']}")
                while self.is_capturing_screen:
                    # Process audio in chunks
                    audio_data = []
                    try:
                        while True:
                            audio_data.append(self.screen_audio_queue.get_nowait())
                    except queue.Empty:
                        if len(audio_data) > 0:
                            # Convert audio data to format suitable for speech recognition
                            audio_data = np.concatenate(audio_data)
                            with tempfile.NamedTemporaryFile(suffix=".wav", delete=True) as temp_file:
                                sf.write(temp_file.name, audio_data, 44100)
                                with sr.AudioFile(temp_file.name) as source:
                                    try:
                                        audio = self.recognizer.record(source)
                                        text = self.recognizer.recognize_google(audio)
                                        if self.is_question(text):
                                            capitalized_text = text[0].upper() + text[1:]
                                            if not capitalized_text.endswith('?'):
                                                capitalized_text += '?'
                                            eel.update_ui(f"Q: {capitalized_text}", "")
                                            self.is_speaking = True
                                            response = self.get_ai_response(capitalized_text)
                                            eel.update_ui("", f"{response}")
                                            self.is_speaking = False
                                    except sr.UnknownValueError:
                                        pass
                                    except Exception as e:
                                        print(f"Error processing screen audio: {str(e)}")
                    time.sleep(0.1)
        except Exception as e:
            print(f"Error capturing screen audio: {str(e)}")
            self.is_capturing_screen = False

    def setup_clients(self):
        if 'openai' in self.api_keys:
            self.openai_client = OpenAI(api_key=self.api_keys['openai'])
        if 'gemini' in self.api_keys:
            genai.configure(api_key=self.api_keys['gemini'])
            self.gemini_client = genai.GenerativeModel('gemini-pro')

    def set_api_key(self, api_key, provider='openai'):
        self.api_keys[provider] = api_key
        self.setup_clients()
        self.save_config()

    def delete_api_key(self, provider='openai'):
        if provider in self.api_keys:
            del self.api_keys[provider]
            if provider == 'openai':
                self.openai_client = None
            elif provider == 'gemini':
                self.gemini_client = None
        self.save_config()

    def has_api_key(self, provider='openai'):
        return provider in self.api_keys

    def set_model(self, model):
        if model in ['gpt', 'gemini']:
            self.current_model = model
            self.save_config()
            return True
        return False

    def get_current_model(self):
        return self.current_model

    def toggle_listening(self):
        if not self.openai_client and not self.gemini_client:
            return False
        self.is_listening = not self.is_listening
        if self.is_listening:
            threading.Thread(target=self.listen_and_process, daemon=True).start()
        return self.is_listening

    def listen_and_process(self):
        cooldown_time = 2  # Cooldown period in seconds
        last_speak_time = 0
        
        while self.is_listening:
            current_time = time.time()
            if not self.is_speaking and not self.audio_playing and (current_time - last_speak_time) > cooldown_time:
                try:
                    with self.mic as source:
                        audio = self.recognizer.listen(source, timeout=5, phrase_time_limit=5)
                    text = self.recognizer.recognize_google(audio)
                    if self.is_question(text):
                        capitalized_text = text[0].upper() + text[1:]
                        if not capitalized_text.endswith('?'):
                            capitalized_text += '?'
                        eel.update_ui(f"Q: {capitalized_text}", "")
                        self.is_speaking = True
                        response = self.get_ai_response(capitalized_text)
                        eel.update_ui("", f"{response}")
                        self.is_speaking = False
                        last_speak_time = time.time()  # Update the last speak time
                except sr.WaitTimeoutError:
                    pass
                except sr.UnknownValueError:
                    pass
                except Exception as e:
                    eel.update_ui(f"An error occurred: {str(e)}", "")
            else:
                time.sleep(0.1)  # Short sleep to prevent busy waiting

    def is_question(self, text):
        # Convert to lowercase for easier matching
        text = text.lower().strip()
        
        # List of question words and phrases
        question_starters = [
            "what", "why", "how", "when", "where", "who", "which",
            "can", "could", "would", "should", "is", "are", "do", "does",
            "am", "was", "were", "have", "has", "had", "will", "shall"
        ]
        
        # Check if the text starts with a question word
        if any(text.startswith(starter) for starter in question_starters):
            return True
        
        # Check for question mark at the end
        if text.endswith('?'):
            return True
        
        # Check for inverted word order (e.g., "Are you...?", "Can we...?")
        if re.match(r'^(are|can|could|do|does|have|has|will|shall|should|would|am|is)\s', text):
            return True
        
        # Check for specific phrases that indicate a question
        question_phrases = [
            "tell me about", "i'd like to know", "can you explain",
            "i was wondering", "do you know", "what about", "how about"
        ]
        if any(phrase in text for phrase in question_phrases):
            return True
        
        # If none of the above conditions are met, it's probably not a question
        return False

    def get_role_prompt(self):
        if self.current_role == "interviewer":
            return self.get_personality_prompt()
        elif self.current_role == "developer":
            return """You are a technical interview expert specializing in software development. Your responses should be:
1. Technically precise but easy to understand
2. Focused on practical coding scenarios
3. Include best practices and common pitfalls

Format your response as:
🔍 Technical Analysis:
[Brief technical explanation]

💻 Code Example (if applicable):
[Simple, relevant code snippet]

⭐ Best Practice:
[One key technical tip]"""
        elif self.current_role == "pitcher":
            return """You are an expert pitch coach for hackathons. Your role is to help participants:
1. Structure their pitch effectively
2. Align with judging criteria
3. Address SDG goals effectively
4. Present technical solutions clearly

Format your response as:
🎯 Pitch Focus:
[Key points to emphasize]

💡 Solution Alignment:
[How it addresses problem/SDGs]

⚡ Impact Highlight:
[Key impact metrics/potential]

🔑 Technical Value:
[Technical innovation/feasibility]"""
        else:
            return """You are a concise interview expert. Your responses should be:
1. Brief but impactful
2. Customized to the specific question
3. Focused on practical advice

Format your response as:
💡 Key Answer Points:
[2-3 bullet points]

🎯 Sample Answer:
[Concise example]

💪 Pro Tip:
[One specific tip]"""

    def build_context(self, question):
        context_parts = [self.get_role_prompt()]
        
        if self.knowledge_base:
            context_parts.append(f"\nRelevant Knowledge Base:\n{self.knowledge_base}")
        
        if self.current_role == "pitcher":
            hackathon_context = f"""
Hackathon Details:
- Name: {self.hackathon_info['name']}
- Problem Statement: {self.hackathon_info['problem_statement']}
- Judging Criteria: {', '.join(f'{c["name"]} ({c["weight"]}%)' for c in self.hackathon_info['criteria'])}
- SDG Goals: {', '.join(self.hackathon_info['sdg_goals'])}
- Additional Context: {self.hackathon_info['additional_context']}"""
            context_parts.append(hackathon_context)
        elif self.current_role == "interviewer":
            context_parts.append(f"""
Context:
- This is a job interview conversation
- Guidelines: {self.job_guidelines}
- Maintain natural conversation flow
- Ask relevant follow-up questions
- Keep responses concise

Question from candidate: {question}""")
        else:
            context_parts.append(f"""
Context:
- Field/Role: {self.job_guidelines}
- Question: {question}

Provide focused, practical guidance.""")
        
        return "\n".join(context_parts)

    def get_ai_response(self, question):
        try:
            base_prompt = self.get_role_prompt()
            context = self.build_context(question)
            
            if self.current_model == "gpt" and self.openai_client:
                return self.get_openai_response(context, question)
            elif self.current_model == "gemini" and self.gemini_client:
                return self.get_gemini_response(context, question)
            
            return "Please configure an AI model and API key first."
        except Exception as e:
            print(f"Error getting AI response: {str(e)}")
            return f"Error getting AI response: {str(e)}"

    def get_gemini_response(self, context, question):
        try:
            # Choose the appropriate Gemini model based on the context
            if self.current_role == "developer" and "code" in question.lower():
                model = genai.GenerativeModel('gemini-pro-code')
            else:
                model = genai.GenerativeModel('gemini-pro')
            
            # Combine context and question
            prompt = f"{context}\n\nQuestion: {question}"
            
            # Generate response
            response = model.generate_content(prompt)
            return response.text
        except Exception as e:
            print(f"Gemini API error: {str(e)}")
            return f"Error with Gemini API: {str(e)}"

    def get_openai_response(self, context, question):
        try:
            # Combine context and question
            prompt = f"{context}\n\nQuestion: {question}"
            
            # Generate response
            response = self.openai_client.generate(prompt, max_tokens=2048)
            return response.choices[0].text
        except Exception as e:
            print(f"OpenAI API error: {str(e)}")
            return f"Error with OpenAI API: {str(e)}"

    def set_hackathon_info(self, info):
        self.hackathon_info.update(info)
        self.save_config()
        return True

    def get_hackathon_info(self):
        return self.hackathon_info

    def set_company_info(self, info):
        self.company_info = info
        self.save_config()
        return True

    def get_company_info(self):
        return self.company_info

    def set_role(self, role):
        self.current_role = role
        self.save_config()
        return True

    def get_current_role(self):
        return self.current_role

    def set_muted(self, muted):
        self.is_muted = muted

    def speaking_ended(self):
        self.is_speaking = False

    def audio_playback_started(self):
        self.audio_playing = True

    def audio_playback_ended(self):
        self.audio_playing = False
        self.is_speaking = False

    def get_current_personality(self):
        return self.current_personality

    def set_personality(self, personality):
        self.current_personality = personality
        self.save_config()
        return True

    def set_language(self, language):
        self.current_language = language
        self.save_config()
        return True

    def get_current_language(self):
        return self.current_language

assistant = AudioAssistant()

@eel.expose
def upload_file(file_data, file_name):
    try:
        file_path = os.path.join(assistant.temp_dir, file_name)
        with open(file_path, 'wb') as f:
            f.write(base64.b64decode(file_data.split(',')[1]))
        
        content = assistant.process_uploaded_file(file_path)
        assistant.knowledge_base += f"\n--- Content from {file_name} ---\n{content}\n"
        
        os.remove(file_path)  # Clean up temporary file
        return True
    except Exception as e:
        print(f"Error processing upload: {str(e)}")
        return False

@eel.expose
def set_company_info(info):
    return assistant.set_company_info(info)

@eel.expose
def get_company_info():
    return assistant.get_company_info()

@eel.expose
def toggle_listening():
    return assistant.toggle_listening()

@eel.expose
def set_api_key_for_provider(api_key, provider):
    try:
        assistant.set_api_key(api_key, provider)
        return True
    except Exception as e:
        print(f"Error saving {provider} API key: {str(e)}")
        return False

@eel.expose
def delete_api_key_for_provider(provider):
    try:
        assistant.delete_api_key(provider)
        return True
    except Exception as e:
        print(f"Error deleting {provider} API key: {str(e)}")
        return False

@eel.expose
def has_api_key_for_provider(provider):
    return assistant.has_api_key(provider)

@eel.expose
def set_ai_model(model):
    return assistant.set_model(model)

@eel.expose
def get_current_model():
    return assistant.get_current_model()

@eel.expose
def set_muted(muted):
    assistant.is_muted = muted

@eel.expose
def speaking_ended():
    assistant.is_speaking = False

@eel.expose
def audio_playback_started():
    assistant.audio_playing = True

@eel.expose
def audio_playback_ended():
    assistant.audio_playing = False
    assistant.is_speaking = False

@eel.expose
def set_job_guidelines(guidelines):
    assistant.set_job_guidelines(guidelines)

@eel.expose
def get_job_guidelines():
    return assistant.get_job_guidelines()

@eel.expose
def start_screen_audio_capture():
    return assistant.start_screen_audio_capture()

@eel.expose
def stop_screen_audio_capture():
    assistant.stop_screen_audio_capture()

@eel.expose
def set_role(role):
    return assistant.set_role(role)

@eel.expose
def get_current_role():
    return assistant.get_current_role()

@eel.expose
def set_personality(personality):
    return assistant.set_personality(personality)

@eel.expose
def get_current_personality():
    return assistant.get_current_personality()

@eel.expose
def set_language(language):
    return assistant.set_language(language)

@eel.expose
def get_current_language():
    return assistant.get_current_language()

@eel.expose
def set_hackathon_info(info):
    return assistant.set_hackathon_info(info)

@eel.expose
def get_hackathon_info():
    return assistant.get_hackathon_info()

if __name__ == '__main__':
    import socket
    
    def try_port(port):
        try:
            sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            sock.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
            sock.bind(('localhost', port))
            sock.close()
            return True
        except:
            return False

    # Try a few different ports
    ports = [8000, 8001, 8002, 8003, 8004, 8005]
    selected_port = None
    
    for port in ports:
        if try_port(port):
            selected_port = port
            break
    
    if selected_port is None:
        print("Error: Could not find a free port. Please try closing your browser and running the application again.")
        exit(1)
    
    print(f"Starting server on port {selected_port}")
    try:
        eel.start('index.html', size=(960, 840), port=selected_port, host='localhost')
    except Exception as e:
        print(f"Failed to start server: {str(e)}")
        print("Try closing your browser completely and run the application again.")
        exit(1)